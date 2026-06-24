#!/usr/bin/env python3
"""Build a <=10 slide PPTX: Conventional vs Wavy 0-degree belt for TBR."""
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

OUT="/home/user/Sri/patent_analysis"
NAVY=RGBColor(0x1f,0x3b,0x5c); TEAL=RGBColor(0x2a,0x9d,0x8f); AMBER=RGBColor(0xe9,0xa2,0x27)
RUST=RGBColor(0xc1,0x49,0x2b); SLATE=RGBColor(0x5b,0x6b,0x7b); WHITE=RGBColor(0xff,0xff,0xff)
LGREY=RGBColor(0xec,0xef,0xf1); DARK=RGBColor(0x22,0x2b,0x33); GREEN=RGBColor(0x3a,0x7d,0x44)

prs=Presentation(); prs.slide_width=IN(13.333); prs.slide_height=IN(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color):
    sp=s.shapes.add_shape(1,IN(x),IN(y),IN(w),IN(h)); sp.fill.solid()
    sp.fill.fore_color.rgb=color; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp_after=4,line=1.04):
    tb=s.shapes.add_textbox(IN(x),IN(y),IN(w),IN(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.vertical_anchor=anchor
    if not runs: return tb
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(sp_after); p.line_spacing=line
        for (t,sz,col,b,it) in para:
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.color.rgb=col; f.bold=b; f.italic=it
            f.name="Calibri"
    return tb
def R(t,sz=14,col=DARK,b=False,it=False): return (t,sz,col,b,it)
def header(s,n,title,kicker=None):
    rect(s,0,0,13.333,1.0,NAVY); rect(s,0,1.0,13.333,0.06,TEAL)
    txt(s,0.45,0.10,11.6,0.85,[[R(title,25,WHITE,True,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        txt(s,0.47,0.62,11.0,0.3,[[R(kicker,11.5,RGBColor(0xbf,0xd6,0xe6),False,True)]])
    # slide number chip
    txt(s,12.5,0.32,0.7,0.4,[[R(f"{n:02d}",16,WHITE,True,False)]],align=PP_ALIGN.RIGHT)
def footer(s,cite):
    rect(s,0,7.18,13.333,0.32,LGREY)
    txt(s,0.45,7.19,12.4,0.3,[[R(cite,8.5,SLATE,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
def pic(s,path,x,y,w):
    img=Image.open(path); ar=img.height/img.width; h=w*ar
    s.shapes.add_picture(path,IN(x),IN(y),IN(w),IN(h)); return h
def bullets(s,x,y,w,h,items,size=13.5,gap=5):
    runs=[]
    for it in items:
        if isinstance(it,tuple): lead,rest=it; runs.append([R("▪  ",size,TEAL,True),R(lead+" ",size,NAVY,True),R(rest,size,DARK)])
        else: runs.append([R("▪  ",size,TEAL,True),R(it,size,DARK)])
    txt(s,x,y,w,h,runs,sp_after=gap,line=1.05)
def card(s,x,y,w,h,title,color,lines,tsize=12.5,bsize=11):
    rect(s,x,y,w,h,WHITE); b=s.shapes[-1]; b.line.color.rgb=color; b.line.width=Pt(1.5)
    rect(s,x,y,w,0.42,color)
    txt(s,x+0.12,y+0.02,w-0.2,0.4,[[R(title,tsize,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    runs=[[R("• ",bsize,color,True),R(l,bsize,DARK)] for l in lines]
    txt(s,x+0.14,y+0.5,w-0.28,h-0.55,runs,sp_after=3,line=1.03)

# ============================================================ SLIDE 1 TITLE
s=slide(); rect(s,0,0,13.333,7.5,NAVY)
rect(s,0,2.55,13.333,0.06,TEAL); rect(s,0,4.30,13.333,0.02,RGBColor(0x35,0x55,0x76))
txt(s,0.9,2.75,11.5,1.4,[
    [R("Zero-Degree Belt Technology for TBR",34,WHITE,True)],
    [R("Conventional  vs  Wavy (Undulating)  Construction",24,RGBColor(0x9f,0xd6,0xe6),False)],
])
txt(s,0.9,1.15,11.5,1.0,[[R("TECHNOLOGY  ·  PROCESS  ·  MATHEMATICS  ·  DIFFERENTIATION",13,TEAL,True)]])
txt(s,0.9,4.55,11.5,1.6,[
 [R("Scope:  Truck & Bus Radial (all-steel)  ·  process & green-to-cure focus  ·  patent basis 2015–2026 + foundational art",13,RGBColor(0xcf,0xdd,0xe8))],
 [R("Covers:  how waviness is induced  ·  governing mathematics  ·  influencing parameters (A/λ, winding tension, ZD layout, cure-lift)  ·  TBR recommendations",13,RGBColor(0xcf,0xdd,0xe8))],
])
txt(s,0.9,6.5,11.5,0.5,[[R("Prepared for internal R&D review  ·  June 2026  ·  Confidential — engineering working document",11,SLATE,False,True)]])

# ============================================================ SLIDE 2 EXEC SUMMARY
s=slide(); header(s,2,"Executive Summary — the core thesis","Why a wavy 0° belt is a genuine lever, not a gimmick")
bullets(s,0.45,1.25,7.55,5.6,[
 ("Problem.","A conventional 0° steel cord is essentially inextensible (working strain to 'taut' ≈ 0.2%). It fights the green→cure crown expansion (typ. 0.5–3%), so a true rigid 0° steel belt is hard to build flat / single-stage."),
 ("Mechanism.","A wavy cord stores length in a wave. It straightens before going taut, giving 1–3% apparent elongation with no strain in the steel — soft during shaping, stiff in service. This is a 'step (bilinear) modulus'."),
 ("Governing law.","Available elongation ε ≈ π²(A/λ)². The industry's wavy A/λ = 1.5–6% window is exactly what covers TBR cure-lift and restores service rigidity — the ranges fall out of the geometry."),
 ("Three constructions compared.","(1) Conventional straight 0°, (2) Wavy/undulating 0°, (3) High-elongation (HE) steel cord — a material route to the same soft→stiff behaviour."),
],size=14,gap=11)
card(s,8.25,1.25,4.65,5.45,"BOTTOM LINE FOR TBR",RUST,[
 "Target A/λ ≈ 3–5%  →  ε ≈ 0.9–2.4%",
 "Apply via servo OSCILLATING head, or",
 "use HE steel cord (drop-in, low gauge)",
 "Closed-loop WINDING TENSION control",
 "Deploy as 0° SHOULDER/EDGE cover to",
 "cut belt-edge shear → protects bead",
 "OWN the white space: amplitude-GRADED",
 "0° (more wave at shoulder, less centre)",
],tsize=13,bsize=12.5)
footer(s,"Math verified numerically (arc-length integral); A/λ ranges & elongation specs from cited patents — see slides 4–6 and references (slide 10).")

# ============================================================ SLIDE 3 LAYOUT / WHY
s=slide(); header(s,3,"Where the 0° belt sits — and how it reaches the bead","TBR all-steel package; 0° as edge cover or full cover")
h=pic(s,f"{OUT}/wavyZD_6_layout.png",4.35,1.30,8.7)
bullets(s,0.4,1.30,3.85,5.5,[
 ("Package.","B1 transition, B2/B3 working belts (~18–22°), B4 protective, plus the 0° circumferential layer."),
 ("Hoop role.","The 0° belt restrains crown radial growth (centrifugal + inflation), keeping the toroid stable at speed."),
 ("The catch.","A stiff straight 0° creates an abrupt stiffness step at its edge → belt-edge / shoulder shear, which is routed down the carcass to the turn-up end (the TBR durability-limiting site)."),
 ("Wavy fix.","A compliant wavy edge softens that step → lower shoulder shear and lower cure-lift stress, without over-stiffening the crown centre (avoids uneven wear)."),
],size=12.5,gap=7)
footer(s,"Belt-cord circumferential elongation spec 1.0–2.5% @ 100–300 N; 17–30 ends/50 mm (crown-reinforcement disclosures). Layout schematic — indicative.")

# ============================================================ SLIDE 4 MATH ELONGATION
s=slide(); header(s,4,"The mathematics of 'free stretch'","How much apparent elongation a wave provides vs its A/λ ratio")
h=pic(s,f"{OUT}/wavyZD_1_math_elongation.png",4.55,1.20,8.5)
card(s,0.4,1.22,3.95,2.05,"GOVERNING FORMULAS",NAVY,[],tsize=12)
txt(s,0.54,1.74,3.7,1.5,[
 [R("Path:  y(x) = A·sin(2πx/λ)",12.5,DARK,True)],
 [R("Exact:  ε = (1/λ)∫₀^λ √(1+(y′)²)dx − 1",12,DARK,False)],
 [R("Small-amplitude:  ε ≈ π²(A/λ)²",13,RUST,True)],
],sp_after=6)
bullets(s,0.4,3.35,3.95,3.6,[
 ("Idea.","The wave stores extra arc-length. As the tyre grows, the wave flattens and the cord 'elongates' geometrically — the steel never strains."),
 ("Verified.","A/λ = 1.5% → 0.22% ; 4% → 1.56% ; 6% → 3.46% (numeric)."),
 ("Design fit.","The 1.5–6% band brackets TBR cure-lift (0.5–3%) and the 1.0–2.5% working-elongation spec."),
],size=12,gap=6)
footer(s,"Ranges: corrugation A/λ 2–15% [US6659147]; wavy/zig-zag A/λ 0.015–0.06 [US8079392, US5054532]. Approximation from arc-length of a sinusoid (standard result).")

# ============================================================ SLIDE 5 BILINEAR
s=slide(); header(s,5,"The 'step / bilinear modulus' — buildable AND durable","Soft through the lift band, stiff in service")
h=pic(s,f"{OUT}/wavyZD_2_load_elongation.png",4.55,1.20,8.5)
bullets(s,0.4,1.25,3.95,5.5,[
 ("Conventional.","Stiff and linear from ε = 0 → very high cord force during shaping; cannot be built flat without special process."),
 ("Wavy.","Bilinear — low modulus until the 'knee' at ε_geo (wave straightened), then full steel modulus."),
 ("HE cord.","Smooth concave low→high modulus from structural elongation in the cord itself."),
 ("Design lever.","Place the knee just past the cure-lift band → service rigidity engages early with no build penalty; lowers belt-edge reaction and controls high-speed standing waves."),
],size=12.3,gap=7)
footer(s,"Soft-stretch / single-stage flat build [US3956546, US3900062, US4050973]; HE steel cord ε_break ≥ 5% [US2016/0152082]; standing-wave cord stress vs speed [Sun et al., Shock & Vibration 2019].")

# ============================================================ SLIDE 6 PROCESS
s=slide(); header(s,6,"Process — four ways to induce the waviness","…and the parameters each route lets you set")
h=pic(s,f"{OUT}/wavyZD_3_process_routes.png",0.5,1.20,9.0)
cards_x=9.75
card(s,cards_x,1.22,3.2,5.6,"ROUTE NOTES",TEAL,[
 "1 Pre-crimp: A=roll depth, λ=tooth",
 "   pitch (3–15 mm; amp 0.1–5 mm)",
 "2 Oscillating head: A=stroke,",
 "   λ=freq÷feed — best control",
 "3 Soft-stretch tape: 'on-end'",
 "   undulations → single-stage build",
 "4 HE cord: structural elong.,",
 "   drop-in, lowest gauge",
 "",
 "Shared: WINDING TENSION ↑",
 "→ uniformity ↑ (−6% LFV pk-pk,",
 "−4% RFV 1st harmonic) but",
 "available ε ↓ & residual tension ↑",
],tsize=12,bsize=10.5)
footer(s,"[US8720175] crimped flat-wire core · [US3956546/US3900062] undulated tape · [US2016/0152082] HE cord · winding-tension uniformity: Procedia Manufacturing (cap-ply cord tension study).")

# ============================================================ SLIDE 7 PARAMETERS
s=slide(); header(s,7,"Influencing parameters — favourability map","Green = favourable; every column framed so 'more = better'")
h=pic(s,f"{OUT}/wavyZD_4_parameters.png",0.5,1.20,8.7)
bullets(s,9.45,1.25,3.5,5.5,[
 ("↑ A/λ.","Big cure-lift headroom, better edge-shear relief — but gauge/weight & uniformity penalty."),
 ("↑ Winding tension.","Uniformity ↑↑, but pulls the wave out (ε ↓) and raises residual hoop tension."),
 ("HE-cord swap.","Broadly favourable — uniformity & lightness gains, simpler line."),
 ("Edge-cover only.","Targets belt-edge shear with lightness; intentionally lowers centre rigidity (good for even wear)."),
],size=11.8,gap=7)
footer(s,"Cells = net favourability (engineering judgement, indicative). Uniformity direction from cap-ply winding-tension study (Procedia Manufacturing).")

# ============================================================ SLIDE 8 COMPARISON
s=slide(); header(s,8,"Head-to-head — conventional vs wavy vs HE cord","Choosing the right construction for the duty")
h=pic(s,f"{OUT}/wavyZD_5_radar.png",6.55,1.20,5.3)
bullets(s,0.4,1.30,4.15,5.4,[
 ("Conventional straight.","Best on simplicity, cost, uniformity, rolling resistance. Weak on cure-lift accommodation and single-stage build."),
 ("Wavy.","Wins cure-lift, flat build, belt-edge/shoulder endurance. Costs gauge/weight, some RR & uniformity, process complexity."),
 ("HE cord.","The balanced middle — most of the wavy benefit with better uniformity and lower gauge, at a special-cord cost."),
 ("Pick by duty.","Long-haul/retread casing → wavy or HE edge cover; high-speed coach → consider full 0°; cost-driven regional → straight + process tweak."),
],size=12.3,gap=8)
footer(s,"Radar scores (1–5, higher = better) are indicative engineering judgement synthesised from the cited art — not measured data.")

# ============================================================ SLIDE 9 RECOMMENDATIONS
s=slide(); header(s,9,"Recommendations & IP white space — TBR","A differentiation playbook")
h=pic(s,f"{OUT}/wavyZD_7_recommendations.png",4.9,1.18,7.35)
bullets(s,0.4,1.25,3.95,4.7,[
 ("Window.","A/λ ≈ 3–5%; ends 17–30/50 mm; knee just past cure-lift."),
 ("Process.","Servo-oscillated application or HE cord; closed-loop tension; consider hybrid (mild wave + HE)."),
 ("Placement.","0° shoulder/edge cover first; full cover only for high-speed."),
 ("Differentiate.","Amplitude-graded 0°; tie to retread / multi-life casing; 6PPD-free coat compound."),
],size=12,gap=6)
card(s,0.4,6.05,12.5,0.92,"NEXT STEPS",NAVY,[
 "FEA + DoE on a real TBR crown (sweep A/λ, winding tension, edge-cover width) to convert indicative ratings into measured claims  ·  Targeted patentability scan on amplitude-graded / position-varying 0° belts before investment.",
],tsize=12,bsize=11.5)
footer(s,"Design window reconciles cited A/λ ranges with TBR cure-lift & working-elongation specs; white-space claim to be confirmed by formal FTO/patentability search.")

# ============================================================ SLIDE 10 REFERENCES
s=slide(); header(s,10,"References & methodology","Primary patents, papers, and basis of estimates")
txt(s,0.45,1.2,6.3,5.6,[
 [R("PATENTS  (process & structure)",13,NAVY,True)],
 [R("US3956546 / US3900062 / US4050973 / US3979536 — 'high soft-stretch' undulated belt tape; single-stage flat-band 0° building (Goodyear).",10.5,DARK)],
 [R("US6659147 — corrugated crown reinforcement; A/λ 2–15% to avoid interfering with shaping during vulcanisation while giving crown rigidity.",10.5,DARK)],
 [R("US8079392 — alternating straight/wavy reinforcement; A/λ 0.015–0.06.",10.5,DARK)],
 [R("US5054532 — wavy / zig-zag cord ply (belt–carcass).",10.5,DARK)],
 [R("US8720175 — crimped flat-wire cord core; pitch 3–15 mm, amplitude 0.1–5 mm.",10.5,DARK)],
 [R("US2016/0152082 A1 — high-elongation steel cord (ε_break ≥ 5%) for 0° / protective layers.",10.5,DARK)],
 [R("US7337817 / US5591284 — circumferential spirally-wound & triangulated steel-cord belts (context).",10.5,DARK)],
],sp_after=5,line=1.03)
txt(s,7.0,1.2,5.9,5.6,[
 [R("PAPERS",13,NAVY,True)],
 [R("Influence of textile cord tension in cap-ply production — Procedia Manufacturing (ScienceDirect): +unwind tension → −6% lateral-force pk-pk, −4% 1st-harmonic radial force.",10.5,DARK)],
 [R("Sun et al. (2019) — Experiment & Analysis of Cord Stress on High-Speed Radial Tire Standing Waves, Shock and Vibration: wavelength & cord stress rise with speed.",10.5,DARK)],
 [R("",6,DARK)],
 [R("METHODOLOGY & LIMITS",13,NAVY,True)],
 [R("• ε(A/λ) curve: exact sinusoid arc-length, numerically verified; small-amplitude form ε ≈ π²(A/λ)².",10.5,DARK)],
 [R("• A/λ ranges, pitch/amplitude, elongation specs: taken from the cited patents.",10.5,DARK)],
 [R("• Bilinear curve shape is schematic but physically grounded (geometry + steel modulus).",10.5,DARK)],
 [R("• Radar scores and the parameter favourability map are indicative engineering judgement, not measured data — for direction-setting, not spec-fixing.",10.5,RUST)],
 [R("• Patent numbers are representative anchors surfaced via Google Patents / Justia / USPTO; verify status & jurisdiction before any FTO use.",10.5,SLATE,False,True)],
],sp_after=5,line=1.03)
footer(s,"Working document for internal R&D · figures generated from first-principles models + cited ranges · confirm all citations against primary sources before external use.")

path=f"{OUT}/Wavy_ZeroDegree_Belt_TBR.pptx"; prs.save(path)
print("SAVED:",path,"| slides:",len(prs.slides._sldIdLst))

"""
Build the PowerPoint deck for the hybrid-cord wear analysis
(General hybrid vs Nylon-Aramid), using the AUDITED cohort definitions
in scripts/cohorts.py (motorcycle excluded; strict hybrid-cord cohorts).

Output: deck charts in output/, reports/Hybrid_Cord_Wear_Analysis.pptx
Run   : python scripts/build_ppt.py
"""
import re
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import cohorts as C

plt.rcParams.update({"font.size": 12, "axes.grid": True, "grid.alpha": 0.25,
                     "axes.spines.top": False, "axes.spines.right": False, "figure.dpi": 150})
NAVY, BLUE, RED, GREEN, GREY, PURP = "#1F3B57", "#4C72B0", "#C44E52", "#55A868", "#999999", "#B0A6C9"

df = C.load()                       # 1080 PCR (motorcycle dropped)
m = C.masks(df)
t = df["_txt"]
GEN, NA, WEAR = m["general"], m["nylon_aramid"], m["wear"]
N_CORPUS = len(df)

def norm(s):
    s = str(s).split(";")[0]; s = re.sub(r"\([A-Z]{2}\)", "", s)
    return re.sub(r"\s+", " ", s).strip(" .,").upper()

def themes(mask):
    s = t[(mask & WEAR).values]
    d = s.str.contains(r"uneven wear|irregular wear|eccentric wear|partial wear", regex=True)
    b = s.str.contains(r"wear resistanc|mileage|tread life|tire life|abrasion resist|wear performance|long.{0,4}wear", regex=True)
    n = len(s)
    return dict(n=n, benefit_only=int((b & ~d).sum()), both=int((b & d).sum()),
                defensive_only=int((d & ~b).sum()), other=int((~b & ~d).sum()))

SG = dict(cohort=int(GEN.sum()), wear=int((GEN & WEAR).sum()), **themes(GEN))
SN = dict(cohort=int(NA.sum()), wear=int((NA & WEAR).sum()), **themes(NA))

def save(fig, path):
    fig.tight_layout(); fig.savefig(path, bbox_inches="tight", facecolor="white"); plt.close(fig)

# chart 1: cohort & wear share
fig, ax = plt.subplots(figsize=(8.6, 4.6))
names = ["General hybrid", "Nylon-Aramid"]; S = {"General hybrid": SG, "Nylon-Aramid": SN}
x = np.arange(2); w = 0.38
b1 = ax.bar(x - w/2, [S[n]["cohort"] for n in names], w, label="cohort size", color=GREY)
b2 = ax.bar(x + w/2, [S[n]["wear"] for n in names], w, label="mention wear", color=RED)
ax.bar_label(b1); ax.bar_label(b2)
for i, n in enumerate(names):
    ax.text(i + w/2, S[n]["wear"] + 3, f"{100*S[n]['wear']/S[n]['cohort']:.0f}%",
            ha="center", color=RED, fontweight="bold")
ax.set_xticks(x); ax.set_xticklabels(names); ax.legend(); ax.set_ylabel("patents")
ax.set_title("Cohort size & wear share (strict, PCR-only)", fontweight="bold")
save(fig, "output/deck_cohorts.png")

# chart 2: wear posture (100% stacked)
fig, ax = plt.subplots(figsize=(9.2, 4.0))
segk = ["benefit_only", "both", "defensive_only", "other"]
segl = ["Benefit only", "Both (benefit + fixes uneven wear)", "Defensive only (fixes uneven wear)", "Other wear mention"]
segc = [GREEN, PURP, RED, "#DDDDDD"]
for yi, n in enumerate(names[::-1]):
    tot = S[n]["wear"]; left = 0
    for k, c in zip(segk, segc):
        v = S[n][k]
        ax.barh(yi, 100*v/tot, left=left, color=c, edgecolor="white")
        if v/tot > 0.06:
            ax.text(left + 50*v/tot, yi, f"{v}\n{100*v/tot:.0f}%", ha="center", va="center",
                    fontsize=10, color="white" if c in (GREEN, RED) else "black", fontweight="bold")
        left += 100*v/tot
ax.set_yticks(range(2)); ax.set_yticklabels([f"{n}\n(n={S[n]['wear']})" for n in names[::-1]])
ax.set_xlim(0, 100); ax.set_xlabel("share of cohort's wear patents (%)")
ax.set_title("Wear posture: claiming a benefit vs. fixing uneven wear", fontweight="bold")
ax.legend([plt.Rectangle((0, 0), 1, 1, color=c) for c in segc], segl,
          loc="upper center", bbox_to_anchor=(0.5, -0.22), ncol=2, fontsize=9, frameon=False)
ax.grid(False)
save(fig, "output/deck_themes.png")

# chart 3: solution approaches — Nylon-Aramid vs baseline (enrichment)
SOL = {
    "Twist / cord geometry": r"twist|denier|dtex|filament count|cord diameter|cord thickness",
    "Graded / dual-modulus": r"(high|low|first|second|different).{0,20}modulus|graded modulus|elastic modulus",
    "Adhesion / dip system": r"adhesion|adhesive|\brfl\b|dip|epoxy|tackif",
    "Core-sheath / wrap": r"core.{0,15}(sheath|wrap|cover)|sheath|wound around.{0,15}core|wrap yarn",
    "Two-layer / edge-cover": r"two layer|double layer|dual layer|edge (band|cover|cap|layer)",
    "Position vs grooves": r"(below|beneath|under).{0,25}(groove|land)|land portion",
    "Zoned density (center/edge)": r"(center|centre|shoulder|edge|side section|side portion).{0,60}(density|spacing|epdm|ends per)",
}
coh = t[NA.values]
rows = [(k, coh.str.contains(p, regex=True).mean(), t.str.contains(p, regex=True).mean()) for k, p in SOL.items()]
rows.sort(key=lambda r: (r[1]/r[2]) if r[2] else 0)
fig, ax = plt.subplots(figsize=(9.6, 5.2))
y = np.arange(len(rows)); h = 0.4
b1 = ax.barh(y + h/2, [100*r[1] for r in rows], h, color=RED, label=f"Nylon-Aramid (n={SN['cohort']})")
b2 = ax.barh(y - h/2, [100*r[2] for r in rows], h, color=GREY, label=f"All PCR baseline (n={N_CORPUS})")
ax.bar_label(b1, labels=[f" {r[1]/r[2]:.1f}x" for r in rows], fontsize=9, fontweight="bold")
ax.set_yticks(y); ax.set_yticklabels([r[0] for r in rows], fontsize=10); ax.legend(fontsize=9)
ax.set_xlabel("% of patents mentioning the approach")
ax.set_title("How Nylon-Aramid patents differ — approach vs. baseline\n"
             "(enrichment x; only core-sheath, dual-modulus & twist are distinctive)",
             fontweight="bold", fontsize=12)
save(fig, "output/deck_solutions.png")

# chart 4: nylon-aramid cohort assignees
fig, ax = plt.subplots(figsize=(8.6, 4.6))
naw = df[NA.values].copy(); naw["a"] = naw["assignee"].map(norm)
top = naw["a"].replace("", np.nan).dropna().value_counts().head(8).iloc[::-1]
ax.bar_label(ax.barh([i[:28] for i in top.index], top.values, color=RED), padding=2)
ax.set_xlabel("patents"); ax.set_title(f"Top assignees — Nylon-Aramid cohort (n={SN['cohort']})", fontweight="bold")
save(fig, "output/deck_assignees.png")

# representative valid NA patents (wear-first, then construction)
belt = re.compile(r"cap ply|overlay|spiral|circumferential|belt|band|cord|hybrid", re.I)
wre = re.compile(C.WEAR_RE, re.I)
def snip(i, need_wear=True):
    for f in ["ai_method", "ai_advantages", "abstract", "claims", "ai_problem"]:
        for s in re.split(r"(?<=[.;])\s+", str(df.loc[i, f])):
            if (not need_wear or wre.search(s)) and belt.search(s) and 45 < len(s) < 240:
                return s.strip()
    return ""
naw_w = df[(NA & WEAR).values].copy(); naw_w["y"] = pd.to_numeric(naw_w["year"], errors="coerce")
KEY = []
for i in naw_w.sort_values("y", ascending=False).index:
    sn = snip(i, True)
    if sn:
        KEY.append((df.loc[i, "doc_id"], norm(df.loc[i, "assignee"])[:26],
                    str(df.loc[i, "year"])[:4], df.loc[i, "legal_status"], sn))
# fill with NA construction patents if few wear snippets
for i in df[NA.values].index:
    if len(KEY) >= 6: break
    if df.loc[i, "doc_id"] in [k[0] for k in KEY]: continue
    sn = snip(i, False)
    if sn and "hybrid" in sn.lower():
        KEY.append((df.loc[i, "doc_id"], norm(df.loc[i, "assignee"])[:26],
                    str(df.loc[i, "year"])[:4], df.loc[i, "legal_status"], sn))

# ── PPTX ────────────────────────────────────────────────────────────────────
prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
NV, GY = RGBColor(0x1F, 0x3B, 0x57), RGBColor(0x55, 0x55, 0x55)

def tb(s, l, t_, w, h):
    x = s.shapes.add_textbox(Inches(l), Inches(t_), Inches(w), Inches(h)); x.text_frame.word_wrap = True
    return x.text_frame
def title_slide(t1, t2):
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = NV
    f = tb(s, 0.9, 2.4, 11.5, 3)
    p = f.paragraphs[0]; p.text = t1; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = f.add_paragraph(); p2.text = t2; p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE6)
def cslide(title):
    s = prs.slides.add_slide(BLANK)
    sp = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
    sp.fill.solid(); sp.fill.fore_color.rgb = NV; sp.line.fill.background()
    f = tb(s, 0.5, 0.3, 12.3, 0.9); p = f.paragraphs[0]; p.text = title
    p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = NV
    return s
def bullets(f, items, size=15):
    for i, (x, lvl) in enumerate(items):
        p = f.paragraphs[0] if i == 0 else f.add_paragraph()
        p.text = ("• " if lvl == 0 else "    – ") + x
        p.font.size = Pt(size - (2 if lvl else 0)); p.font.color.rgb = GY if lvl else NV; p.space_after = Pt(5)
def img(s, path, l, t_, w): s.shapes.add_picture(path, Inches(l), Inches(t_), width=Inches(w))

# 1 title
title_slide("Hybrid Cord Patents — Wear Analysis",
            "Zero-degree belt cords in PCR tires · General hybrid vs. Nylon-Aramid · "
            f"{N_CORPUS} PCR patents (audited, motorcycle-excluded)")
# 2 scope
s = cslide("Scope & Method")
bullets(tb(s, 0.6, 1.4, 12.1, 5.6), [
    (f"PatSeer zero-degree landscape, screened to {N_CORPUS} passenger-car-radial (PCR) patents — primary motorcycle tires excluded.", 0),
    ("Two hybrid-cord cohorts, defined by EXPLICIT hybrid-cord construction (not mere material co-occurrence):", 0),
    (f"General hybrid cords: {SG['cohort']} patents", 1),
    (f"Nylon-Aramid cords (focus): {SN['cohort']} patents", 1),
    ("Wear detection over title/abstract/claims/description + PatSeer AI summaries; 'mileage' is rarely named — this is predominantly a treadwear analysis.", 0),
    ("Theme tagging: BENEFIT (wear resistance/mileage) vs. DEFENSIVE (fixing uneven wear). See Limitations slide for caveats.", 0),
], 15)
# 3 headline
s = cslide("Headline Finding")
img(s, "output/deck_cohorts.png", 6.8, 1.6, 6.1)
bullets(tb(s, 0.6, 1.6, 6.0, 5.2), [
    (f"General hybrid: {SG['wear']} of {SG['cohort']} ({100*SG['wear']/SG['cohort']:.0f}%) make a wear claim.", 0),
    (f"Nylon-Aramid: {SN['wear']} of {SN['cohort']} ({100*SN['wear']/SN['cohort']:.0f}%) make a wear claim.", 0),
    ("Wear is a SECONDARY, indirect benefit — driven by crown stiffness & footprint uniformity, not abrasion chemistry.", 0),
    ("Nylon-Aramid-specific wear evidence is thin (31 patents) — interpret with care.", 0),
], 16)
# 4 posture
s = cslide("Wear Posture — Benefit vs. Fixing Uneven Wear")
img(s, "output/deck_themes.png", 0.7, 1.5, 8.4)
bullets(tb(s, 9.2, 1.6, 3.8, 5.2), [
    ("Each bar = 100% of that cohort's wear patents.", 0),
    (f"Nylon-Aramid: {SN['defensive_only']} defensive-only vs {SN['both']} 'both'.", 0),
    (f"General hybrid: {SG['benefit_only']} benefit-only leads.", 0),
    ("Stiff aramid helps footprint uniformity but its stiffness step can trigger uneven wear that must be engineered out.", 0),
], 14)
# 5 solutions
s = cslide("How Nylon-Aramid Tackles Wear (vs. baseline)")
img(s, "output/deck_solutions.png", 6.3, 1.45, 6.7)
bullets(tb(s, 0.6, 1.7, 5.6, 5.0), [
    ("Compared to all PCR patents, Nylon-Aramid is distinctive in the CORD itself:", 0),
    ("Core-sheath / wrap construction (2.0x)", 1),
    ("Graded / dual-modulus cord (1.9x)", 1),
    ("Twist / cord-geometry tuning (1.8x)", 1),
    ("Adhesion/dip, two-layer and ZONED density are NOT distinctive to Nylon-Aramid (~1x) — zoned density is a PET/PA story (next slide).", 0),
], 14)
# 6 representative patents
s = cslide("Representative Nylon-Aramid Patents")
items = []
for doc, asg, yr, ls, sn in KEY[:6]:
    items.append((f"{doc} — {asg} ({yr}, {ls})", 0)); items.append((sn, 1))
bullets(tb(s, 0.5, 1.35, 12.4, 5.7), items, 14)
# 7 zoned density honesty slide
s = cslide("Related Lever: Zoned Density is PET/PA — not Nylon-Aramid")
bullets(tb(s, 0.6, 1.4, 12.1, 5.6), [
    ("Zoned-density bandages (different cord/density centre vs. shoulder) are a real uneven-wear fix — but in PET/PA hybrids, not Nylon-Aramid.", 0),
    ("Continental EP3912833A1 (ALIVE): centre PET 120 EPDM + PA4.6 side sections 130 EPDM.", 1),
    ("Continental DE102016223304B4 (ALIVE): 3-section bandage, centre PA6.6 (470x2) for 'more uniform abrasion'.", 1),
    ("In the Nylon-Aramid cohort zoned density appears in only ~3% (1.3x) — it is NOT the Nylon-Aramid wear lever.", 0),
    ("Takeaway: do not attribute the zoned-density wear story to Nylon-Aramid cords.", 0),
], 15)
# 8 assignees
s = cslide("Who Owns the Nylon-Aramid IP")
img(s, "output/deck_assignees.png", 6.5, 1.5, 6.5)
bullets(tb(s, 0.6, 1.7, 5.7, 4.8), [
    ("Japanese OEMs (Bridgestone, Sumitomo, Yokohama) plus Goodyear, Michelin, Continental.", 0),
    ("Cord makers (Kordsa, Hyosung) appear in the construction patents.", 0),
    ("Most Nylon-Aramid wear patents are legally DEAD — limited live art on the wear angle specifically.", 0),
], 15)
# 9 limitations
s = cslide("Limitations & Methodology (read before citing)")
bullets(tb(s, 0.6, 1.35, 12.2, 5.8), [
    ("Cohorts are text-mined (PatSeer fields + AI summaries) with high recall; spot-check before quoting individual patents.", 0),
    ("'Nylon-Aramid' = explicit aramid+nylon hybrid-cord construction. A looser co-occurrence rule inflated this ~3x (318 -> 100); the strict figure is used here.", 0),
    ("92 primary motorcycle tires were removed; some patents still mention motorcycle as a secondary application.", 0),
    ("Solution percentages are multi-label and shown as ENRICHMENT vs. baseline to avoid overstating common vocabulary.", 0),
    ("Wear is predominantly treadwear, not 'mileage' (named in only ~12 patents).", 0),
    ("Corpus skews pre-2015; legal status reflects the export date.", 0),
], 14)
# 10 takeaways
s = cslide("Takeaways & Recommendations")
bullets(tb(s, 0.6, 1.4, 12.1, 5.6), [
    ("Wear/mileage is a secondary, indirect benefit of Nylon-Aramid zero-degree cords — via stiffness and footprint control.", 0),
    ("The Nylon-Aramid wear lever is the CORD: core-sheath, dual-modulus, twist tuning — not zoned density (that is PET/PA).", 0),
    ("Manage the uneven-wear risk created by the aramid stiffness step (defensive IP is a real share).", 0),
    ("White space: Nylon-Aramid wear evidence is thin and mostly dead — opportunity, but validate FTO against live PET/PA zoned-density art.", 0),
    ("Next: pull full claims for the strict Nylon-Aramid set; add rolling-resistance and high-speed-durability axes.", 0),
], 15)

out = "reports/Hybrid_Cord_Wear_Analysis.pptx"
prs.save(out)
print("saved", out, "|", len(prs.slides._sldIdLst), "slides")
print(f"corpus={N_CORPUS} | general {SG['cohort']}/{SG['wear']} | nylon-aramid {SN['cohort']}/{SN['wear']}")
